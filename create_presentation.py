"""
create_presentation.py
Generates Churn_Prediction_Presentation.pptx — a dark, modern pitch deck
based on Churn_Prediction_Report_Revised.docx.
Run with: python3 create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUTPUT_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "Churn_Prediction_Presentation.pptx",
)

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
CARD      = RGBColor(0x16, 0x2A, 0x40)   # slightly lighter navy for cards
CYAN      = RGBColor(0x00, 0xC6, 0xFF)   # electric cyan (primary accent)
CORAL     = RGBColor(0xFF, 0x6B, 0x6B)   # coral (danger / highlight)
MINT      = RGBColor(0x00, 0xE5, 0x96)   # mint green (positive metric)
GOLD      = RGBColor(0xFF, 0xD7, 0x00)   # gold (numbers / KPIs)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xB0, 0xC4, 0xD8)   # light blue-gray for body text
DGRAY     = RGBColor(0x44, 0x60, 0x77)   # subtle divider / muted text

# ---------------------------------------------------------------------------
# Fonts  (Helvetica Neue for all text; Impact only for hero stat numbers)
# ---------------------------------------------------------------------------
F_TITLE  = "Helvetica Neue"   # section titles, labels
F_BODY   = "Helvetica Neue"   # body / bullet text
F_HERO   = "Impact"           # giant KPI numbers (93, $6 510, etc.)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, alpha=None):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, font=None):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font or F_TITLE
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def hero(slide, number, label, source, l, t, w, num_color=GOLD):
    """Hero KPI block: big Impact number + Helvetica label + tiny source."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(2.5))
    tf  = txb.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = number
    r1.font.name  = F_HERO
    r1.font.size  = Pt(64)
    r1.font.color.rgb = num_color

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = label
    r2.font.name  = F_BODY
    r2.font.size  = Pt(15)
    r2.font.color.rgb = LGRAY

    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = source
    r3.font.name  = F_BODY
    r3.font.size  = Pt(10)
    r3.font.color.rgb = DGRAY
    return txb


def accent_bar(slide, color=CYAN, height=0.06):
    """Thin horizontal accent bar across the top."""
    rect(slide, 0, 0, 13.33, height, color)


def slide_label(slide, text, color=CYAN):
    """Small all-caps category label in the top-left."""
    txt(slide, text.upper(), 0.45, 0.18, 6, 0.35,
        size=11, bold=True, color=color)


def dot(slide, l, t, size=0.12, color=CYAN):
    """Small circle accent."""
    shape = slide.shapes.add_shape(
        9,  # OVAL
        Inches(l), Inches(t), Inches(size), Inches(size),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = blank_slide(prs)
    bg(s)

    # Left cyan bar
    rect(s, 0, 0, 0.08, 7.5, CYAN)

    # Large diagonal decorative block top-right
    rect(s, 10.5, 0, 2.83, 2.2, CARD)
    rect(s, 11.2, 0, 2.13, 1.4, RGBColor(0x00, 0x7A, 0xCC))

    # Tag line chip
    rect(s, 0.45, 0.85, 3.2, 0.42, CYAN)
    txt(s, "AI-Powered Retention Intelligence", 0.55, 0.88, 3.1, 0.38,
        size=13, bold=True, color=BG, align=PP_ALIGN.LEFT)

    # Main title
    txt(s, "Customer Churn\nPrediction", 0.45, 1.5, 9, 2.4,
        size=72, bold=True, color=WHITE)

    # Subtitle
    txt(s, "From CSV to Executive Brief — Automatically", 0.45, 4.05, 9, 0.6,
        size=24, bold=False, color=LGRAY)

    # Divider
    rect(s, 0.45, 4.78, 4.5, 0.04, CYAN)

    # Meta
    txt(s, "Multi-Agent · LangGraph · SHAP · Bayesian Optimisation",
        0.45, 4.96, 9, 0.45, size=14, color=DGRAY)
    txt(s, "Final Project Report  ·  April 2026  ·  Columbia University",
        0.45, 5.45, 9, 0.45, size=14, color=DGRAY)

    # Bottom-right decorative dots
    for i, c in enumerate([CYAN, MINT, CORAL]):
        dot(s, 12.3 + i * 0.3, 6.9, 0.18, c)


def slide_problem(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, CORAL)
    slide_label(s, "The Problem", CORAL)

    txt(s, "Churn is Expensive.\nAnd Preventable.", 0.45, 0.55, 9, 1.8,
        size=52, bold=True, color=WHITE)

    # Three hero stat cards
    stats = [
        ("$440B+",
         "lost annually to subscription\nchurn & payment failures",
         "Recurly, 2024 State of Subscriptions",
         CORAL),
        ("5% more\nretention",
         "= up to 95% profit increase\nacross industries",
         "Reichheld & Bain, The Loyalty Effect (HBR)",
         GOLD),
        ("17%",
         "annual churn rate at Peacock;\n22% industry avg in telecom",
         "Parks Associates / Antenna, 2024",
         CYAN),
    ]
    for i, (num, label, source, color) in enumerate(stats):
        lx = 0.45 + i * 4.27
        rect(s, lx, 2.55, 4.0, 2.8, CARD)
        rect(s, lx, 2.55, 4.0, 0.06, color)
        txt(s, num,    lx+0.2, 2.68, 3.65, 1.1,  size=38, bold=True,  color=color, font=F_HERO)
        txt(s, label,  lx+0.2, 3.7,  3.65, 0.85, size=13, bold=False, color=WHITE)
        txt(s, source, lx+0.2, 5.15, 3.65, 0.4,  size=10, bold=False, color=DGRAY)

    txt(s,
        "Most organisations know this — yet retention teams still can't act on it because "
        "turning raw data into an actionable list requires weeks of specialist work.",
        0.45, 5.75, 12.4, 0.7, size=15, color=LGRAY)


def slide_old_way(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, DGRAY)
    slide_label(s, "Status Quo", LGRAY)

    txt(s, "The Data Science Bottleneck", 0.45, 0.5, 10, 0.75,
        size=44, bold=True, color=WHITE)

    # Two large research stats across the top
    rect(s, 0.45, 1.45, 5.9, 2.05, CARD)
    rect(s, 0.45, 1.45, 5.9, 0.06, CORAL)
    txt(s, "60%", 0.65, 1.58, 5.5, 1.0, size=62, bold=True, color=CORAL, font=F_HERO)
    txt(s, "of data professionals' time is spent\ngetting to insight — not doing it",
        0.65, 2.48, 5.5, 0.65, size=13, color=WHITE)
    txt(s, "37% searching for data · 36% preparing it",
        0.65, 3.08, 5.5, 0.35, size=11, color=DGRAY)
    txt(s, "Alteryx survey, 2018 (n = 1,000+ data professionals)",
        0.65, 3.3, 5.5, 0.25, size=10, color=DGRAY)

    rect(s, 6.55, 1.45, 6.35, 2.05, CARD)
    rect(s, 6.55, 1.45, 6.35, 0.06, GOLD)
    txt(s, "62%", 6.75, 1.58, 5.9, 1.0, size=62, bold=True, color=GOLD, font=F_HERO)
    txt(s, "of analysts depend on other teams\nto complete analytics steps",
        6.75, 2.48, 5.9, 0.65, size=13, color=WHITE)
    txt(s, "KNIME: Remove Bottlenecks with Self-Service Analytics",
        6.75, 3.3, 5.9, 0.25, size=10, color=DGRAY)

    # Three week boxes (practitioner experience, clearly labelled as such)
    weeks = [
        ("WEEK 1", "Export data from CRM.\nSubmit ticket to data science.\nWait in the queue.", CORAL),
        ("WEEK 2", "DS team cleans data,\ntrains a model, returns\na spreadsheet of scores.", GOLD),
        ("WEEK 3", "Analyst gets scores —\nno 'why', no scenario tool,\nno segment narrative.", LGRAY),
    ]
    for i, (week, body_text, color) in enumerate(weeks):
        lx = 0.45 + i * 4.3
        rect(s, lx, 3.8, 4.0, 2.45, CARD)
        rect(s, lx, 3.8, 4.0, 0.5, color)
        txt(s, week, lx+0.18, 3.84, 3.7, 0.42, size=16, bold=True, color=BG)
        txt(s, body_text, lx+0.18, 4.42, 3.65, 1.65, size=14, color=WHITE)

    txt(s, "Typical practitioner experience — analyst wait time measured in weeks, not hours.",
        0.45, 6.45, 12.4, 0.42, size=13, bold=True, color=CORAL)


def slide_solution(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, MINT)
    slide_label(s, "The Solution", MINT)

    txt(s, "One Upload.\nEverything Else is Automatic.", 0.45, 0.5, 10, 1.8,
        size=48, bold=True, color=WHITE)

    benefits = [
        (CYAN,  "Any Dataset, Zero Setup",
                "Auto-detects schema, columns & labels from any CSV. No config."),
        (MINT,  "Under 2 Minutes",
                "Full pipeline runs in the background. Complete analysis before your next meeting."),
        (GOLD,  "Plain-English Segments",
                "Named customer risk groups with recommended retention plays. No data science degree required."),
        (CORAL, "AI-Written Executive Summary",
                "GPT-4.1 drafts a presentation-ready brief with KPIs, actions & driver narratives."),
        (RGBColor(0xC0, 0x7A, 0xFF), "Interactive Scenarios",
                "Adjust contact budget. Ask 'why is this customer flagged?' All in real time."),
    ]
    for i, (color, title, desc) in enumerate(benefits):
        col = i % 3
        row = i // 3
        lx = 0.45 + col * 4.27
        ty = 2.55 + row * 2.1
        rect(s, lx, ty, 4.0, 1.85, CARD)
        rect(s, lx, ty, 0.06, 1.85, color)
        txt(s, title, lx + 0.2, ty + 0.12, 3.7, 0.5, size=15, bold=True, color=color)
        txt(s, desc,  lx + 0.2, ty + 0.6,  3.7, 1.1, size=13, color=LGRAY)


def slide_pipeline(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)
    slide_label(s, "How It Works")

    txt(s, "10-Stage Automated Pipeline", 0.45, 0.5, 9, 0.7,
        size=44, bold=True, color=WHITE)

    stages = [
        ("01", "Schema\nDetection",   CYAN),
        ("02", "Horizon\nDefinition", CYAN),
        ("03", "Imbalance\nAgent",    CYAN),
        ("04", "Missing\nValues",     CYAN),
        ("05", "Data\nCleaning",      CYAN),
        ("06", "Model\nTraining",     MINT),
        ("07", "SHAP\nExplainer",     MINT),
        ("08", "Business\nAggregates",GOLD),
        ("09", "Segment\nDiscovery",  GOLD),
        ("10", "Insight\nGeneration", CORAL),
    ]

    box_w = 1.17
    box_h = 1.55
    gap   = 0.07
    start_l = 0.28

    for i, (num, label, color) in enumerate(stages):
        lx = start_l + i * (box_w + gap)
        rect(s, lx, 1.55, box_w, box_h, CARD)
        rect(s, lx, 1.55, box_w, 0.32, color)
        txt(s, num,   lx + 0.08, 1.58, box_w - 0.1, 0.3, size=13, bold=True, color=BG)
        txt(s, label, lx + 0.08, 1.96, box_w - 0.1, 1.0, size=12, color=WHITE)

        # Arrow between boxes
        if i < len(stages) - 1:
            ax = lx + box_w + 0.01
            txt(s, "›", ax, 2.05, 0.1, 0.5, size=16, bold=True, color=DGRAY)

    # Phase labels below
    rect(s, 0.28, 3.25, 5 * (box_w + gap) - gap, 0.3, RGBColor(0x0A, 0x14, 0x1F))
    txt(s, "Offline Pipeline — Data & Prep", 0.35, 3.27, 5.8, 0.28,
        size=11, color=CYAN)

    rect(s, 0.28 + 5 * (box_w + gap), 3.25, 2 * (box_w + gap) - gap, 0.3,
         RGBColor(0x0A, 0x14, 0x1F))
    txt(s, "Modelling", 0.28 + 5 * (box_w + gap) + 0.07, 3.27, 2.5, 0.28,
        size=11, color=MINT)

    rect(s, 0.28 + 7 * (box_w + gap), 3.25, 3 * (box_w + gap) - gap, 0.3,
         RGBColor(0x0A, 0x14, 0x1F))
    txt(s, "Business Intelligence", 0.28 + 7 * (box_w + gap) + 0.07, 3.27, 3.5, 0.28,
        size=11, color=GOLD)

    # Online layer callout
    rect(s, 0.28, 3.8, 12.8, 2.95, RGBColor(0x0E, 0x22, 0x35))
    rect(s, 0.28, 3.8, 12.8, 0.06, DGRAY)
    txt(s, "ONLINE AGENT LAYER  —  Available after pipeline completes",
        0.5, 3.83, 10, 0.38, size=12, bold=True, color=DGRAY)

    agents = [
        (CYAN,  "Chart Agent",       "On-demand visualisations"),
        (GOLD,  "Simulation Agent",  "Real-time budget & threshold optimisation"),
        (MINT,  "Insight Agent",     "Grounded Q&A on model results"),
        (CORAL, "Chat Agent",        "Open-ended data exploration"),
    ]
    for i, (color, name, desc) in enumerate(agents):
        lx = 0.5 + i * 3.2
        rect(s, lx, 4.3, 3.0, 1.3, CARD)
        rect(s, lx, 4.3, 3.0, 0.06, color)
        txt(s, name, lx + 0.15, 4.42, 2.75, 0.42, size=14, bold=True, color=color)
        txt(s, desc, lx + 0.15, 4.88, 2.75, 0.65, size=12, color=LGRAY)


def slide_model_results(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, MINT)
    slide_label(s, "Results — Model Performance", MINT)

    txt(s, "XGBoost Wins.\nOn Accuracy & Business Value.", 0.45, 0.5, 9, 1.8,
        size=46, bold=True, color=WHITE)

    models = [
        ("XGBoost",           "0.821", "0.344", "$6,510", True),
        ("LightGBM",          "0.814", "0.331", "—",      False),
        ("Gradient Boosting", "0.809", "0.325", "—",      False),
        ("Random Forest",     "0.791", "0.298", "—",      False),
        ("Logistic Regression","0.763","0.271", "—",      False),
    ]

    headers = ["Model", "ROC-AUC", "PR-AUC", "Expected Profit"]
    col_x   = [0.45, 4.5, 7.0, 9.4]
    col_w   = [3.8,  2.3, 2.1, 3.5]

    # Header row
    rect(s, 0.45, 2.4, 12.5, 0.48, CYAN)
    for j, h in enumerate(headers):
        txt(s, h, col_x[j] + 0.1, 2.44, col_w[j], 0.4,
            size=14, bold=True, color=BG,
            align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    for i, (name, roc, pr, profit, is_best) in enumerate(models):
        ty = 2.92 + i * 0.72
        row_color = RGBColor(0x16, 0x38, 0x52) if is_best else CARD
        rect(s, 0.45, ty, 12.5, 0.68, row_color)
        if is_best:
            rect(s, 0.45, ty, 0.06, 0.68, MINT)

        vals = [name, roc, pr, profit]
        vcols = [MINT if is_best else WHITE, LGRAY, LGRAY,
                 GOLD if is_best else LGRAY]
        for j, (v, vc) in enumerate(zip(vals, vcols)):
            txt(s, v, col_x[j] + 0.1, ty + 0.1, col_w[j], 0.5,
                size=15 if is_best else 14,
                bold=is_best and j == 0,
                color=vc,
                align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    txt(s, "Selected on combined ROC-AUC + expected retention profit",
        0.45, 6.7, 12, 0.4, size=13, color=DGRAY)


def slide_time(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, GOLD)
    slide_label(s, "Results — Speed", GOLD)

    txt(s, "95% Faster.\nSame Output.", 0.45, 0.5, 9, 1.6, size=56, bold=True, color=WHITE)

    # Before box
    rect(s, 0.45, 2.35, 5.8, 3.8, CARD)
    rect(s, 0.45, 2.35, 5.8, 0.55, RGBColor(0x55, 0x22, 0x22))
    txt(s, "BEFORE", 0.65, 2.4, 5.4, 0.45, size=16, bold=True, color=CORAL)
    items = [
        "Week 1 — Submit data science ticket",
        "Week 2 — Model training & scoring",
        "Week 3 — Receive spreadsheet, no explanations",
        "2–3 weeks total",
    ]
    for i, item in enumerate(items):
        txt(s, ("⚠  " if i == 3 else "·  ") + item,
            0.65, 3.05 + i * 0.68, 5.3, 0.6,
            size=15 if i < 3 else 17,
            bold=(i == 3),
            color=CORAL if i == 3 else LGRAY)

    # Arrow
    txt(s, "→", 6.4, 3.9, 0.9, 0.9, size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # After box
    rect(s, 7.45, 2.35, 5.4, 3.8, RGBColor(0x0A, 0x28, 0x1E))
    rect(s, 7.45, 2.35, 5.4, 0.55, RGBColor(0x08, 0x44, 0x2C))
    txt(s, "AFTER", 7.65, 2.4, 5.0, 0.45, size=16, bold=True, color=MINT)

    txt(s, "93", 7.65, 2.98, 5.0, 1.4, size=100, bold=True, color=MINT)
    txt(s, "seconds", 9.95, 3.4, 2.6, 0.6, size=26, bold=False, color=MINT)
    txt(s, "Trained model  ·  SHAP explanations\nNamed segments  ·  Executive Summary\nReady to present",
        7.65, 4.55, 5.0, 1.4, size=14, color=LGRAY)


def slide_money(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, GOLD)
    slide_label(s, "Results — Business Impact", GOLD)

    txt(s, "The Numbers That Matter.", 0.45, 0.5, 10, 0.8, size=46, bold=True, color=WHITE)

    kpis = [
        (GOLD,  "$6,510",
         "Projected retention profit at optimal threshold",
         "This study — XGBoost on 10k-customer test set"),
        (MINT,  "10–15%",
         "Churn reduction from analytics-driven retention\nin telecom (McKinsey)",
         "McKinsey: Reducing Churn in Telecom, Dec 2017"),
        (CYAN,  "10–12%",
         "Recommended outreach pool at threshold 0.168\n(vs. naive 50% default)",
         "This study — business-value threshold optimisation"),
        (CORAL, "Only 41%",
         "of companies can explain their AI model decisions\n— the explainability gap this tool closes",
         "IBM Global AI Adoption Index, Jan 2024"),
    ]

    for i, (color, number, label, source) in enumerate(kpis):
        col = i % 2
        row = i // 2
        lx = 0.45 + col * 6.45
        ty = 1.58 + row * 2.72
        rect(s, lx, ty, 6.15, 2.5, CARD)
        rect(s, lx, ty, 6.15, 0.06, color)
        txt(s, number, lx+0.25, ty+0.12, 5.7, 1.05, size=58, bold=True, color=color, font=F_HERO)
        txt(s, label,  lx+0.25, ty+1.22, 5.7, 0.85, size=14, color=WHITE)
        txt(s, source, lx+0.25, ty+2.14, 5.7, 0.32, size=10, color=DGRAY)


def slide_features(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, CYAN)
    slide_label(s, "Results — What's Driving Churn")

    txt(s, "SHAP: Why Each Customer Was Flagged", 0.45, 0.5, 9, 0.7,
        size=40, bold=True, color=WHITE)
    txt(s, "SHapley values attribute each prediction to individual features. "
           "SHAP & LIME are the two most widely used XAI methods for tabular data "
           "(Salih et al., Advanced Intelligent Systems, 2025).",
        0.45, 1.28, 12.4, 0.55, size=14, color=LGRAY)

    features = [
        ("csat_score",       0.60, CORAL,
         "Dominant signal — consistent with McKinsey telecom research on satisfaction as #1 churn driver.",
         "This study (SHAP on XGBoost, test set)"),
        ("tenure_months",    0.38, GOLD,
         "New customers churn at higher rates. SaaS median annual revenue churn 12.5% — skewed to early tenure.",
         "Lighter Capital 2025 B2B SaaS Benchmarks"),
        ("payment_failures", 0.32, CYAN,
         "Recurly 2024: subscription businesses lose $440B+ annually to payment failures — a leading involuntary churn driver.",
         "Recurly, 2024 State of Subscriptions"),
        ("monthly_logins",   0.30, MINT,
         "39% of US consumers cancelled a streaming service in 6 months; declining engagement precedes cancellation.",
         "Deloitte / Variety Digital Trends, 2023"),
    ]

    bar_max_w = 6.8
    for i, (feat, shap, color, insight, source) in enumerate(features):
        ty = 2.05 + i * 1.2
        bar_w = bar_max_w * (shap / 0.65)

        txt(s, feat,              0.45, ty,       3.05, 0.48, size=16, bold=True, color=color)
        txt(s, f"SHAP ~{shap:.2f}",0.45, ty+0.46, 3.05, 0.38, size=12, color=DGRAY)
        rect(s, 3.65, ty+0.1, bar_max_w, 0.52, RGBColor(0x16, 0x2A, 0x40))
        rect(s, 3.65, ty+0.1, bar_w,     0.52, color)
        txt(s, insight, 3.75, ty+0.66, bar_max_w-0.05, 0.38, size=11, color=LGRAY)
        txt(s, source,  3.75, ty+1.0,  bar_max_w-0.05, 0.26, size=10, color=DGRAY)

    txt(s, "Actionable: CSAT drop = earliest visible churn signal. Monitor it weekly.",
        0.45, 6.88, 12.4, 0.42, size=13, bold=True, color=CORAL)


def slide_conclusion(prs):
    s = blank_slide(prs)
    bg(s)

    # Full-width top block
    rect(s, 0, 0, 13.33, 3.5, RGBColor(0x05, 0x10, 0x1A))
    rect(s, 0, 0, 13.33, 0.08, CYAN)

    txt(s, "The barrier to AI-powered\nretention management is not\nmodel accuracy.",
        0.6, 0.25, 12, 2.9, size=52, bold=True, color=WHITE)

    txt(s, "It is", 0.6, 3.1, 1.5, 0.6, size=30, bold=False, color=LGRAY)
    txt(s, "accessibility.", 1.9, 3.1, 5, 0.65, size=40, bold=True, color=CYAN)

    txt(s,
        "Gradient boosting has been solving churn prediction for a decade. "
        "What's missing is a tool that puts those predictions in the hands of the people who can act on them — "
        "without a data science team in the loop.",
        0.6, 3.9, 12.0, 1.1, size=16, color=LGRAY)

    # Three closing points
    closes = [
        (MINT,  "Upload any CSV →"),
        (GOLD,  "Get a complete brief in 93 seconds →"),
        (CORAL, "Walk into the meeting ready"),
    ]
    for i, (color, line) in enumerate(closes):
        txt(s, line, 0.6 + i * 4.25, 5.2, 4.1, 0.6,
            size=15, bold=True, color=color)

    # Bottom bar
    rect(s, 0, 7.1, 13.33, 0.4, RGBColor(0x08, 0x14, 0x22))
    txt(s, "Columbia University  ·  AI Churn Prediction Tool  ·  April 2026",
        0, 7.14, 13.33, 0.35, size=12, color=DGRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Assemble
# ---------------------------------------------------------------------------
prs = new_prs()

slide_title(prs)
slide_problem(prs)
slide_old_way(prs)
slide_solution(prs)
slide_pipeline(prs)
slide_model_results(prs)
slide_time(prs)
slide_money(prs)
slide_features(prs)
slide_conclusion(prs)

prs.save(OUTPUT_PATH)
print(f"Saved → {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
